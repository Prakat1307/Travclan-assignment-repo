"""
Hotel Booking Analysis – PowerPoint Report Generator
=====================================================
Deliverable 3: Summary Report (PowerPoint)
Covers: Observed patterns, Root causes, Actionable recommendations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, sys

# Ensure UTF-8 output on Windows
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

OUTPUT_DIR = r"e:\assignment\output"
os.makedirs(OUTPUT_DIR, exist_ok=True)
REPORT_PATH = os.path.join(OUTPUT_DIR, "Hotel_Booking_Analysis_Report.pptx")
VIZ_PATH    = os.path.join(OUTPUT_DIR, "visualizations.png")

# ─────────────────────────────────────────────────────────────────────────────
# COLOURS
# ─────────────────────────────────────────────────────────────────────────────
C_INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
C_CYAN    = RGBColor(0x06, 0xB6, 0xD4)
C_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
C_RED     = RGBColor(0xEF, 0x44, 0x44)
C_GREEN   = RGBColor(0x10, 0xB9, 0x81)
C_DARK    = RGBColor(0x1E, 0x29, 0x3B)
C_SLATE   = RGBColor(0x47, 0x55, 0x69)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)
C_INDIGO2 = RGBColor(0x31, 0x2E, 0x81)   # dark indigo for header bg

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, alpha_needed=False):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h, font_size=12, bold=False,
                 color=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_slide_header(slide, title, subtitle=None,
                     bg_color=C_INDIGO2, accent_color=C_CYAN):
    """Full-width header bar."""
    add_rect(slide, 0, 0, 13.33, 1.3, bg_color)
    add_rect(slide, 0, 1.3, 13.33, 0.06, accent_color)
    add_text_box(slide, title, 0.35, 0.12, 12.6, 0.75,
                 font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.38, 0.82, 12.5, 0.45,
                     font_size=12, color=C_CYAN, align=PP_ALIGN.LEFT)


def add_kpi_card(slide, l, t, w, h, label, value, sub="", bg=C_INDIGO):
    add_rect(slide, l, t, w, h, bg)
    add_text_box(slide, value, l+0.1, t+0.08, w-0.2, 0.5,
                 font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l+0.1, t+0.54, w-0.2, 0.35,
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if sub:
        add_text_box(slide, sub, l+0.1, t+0.85, w-0.2, 0.28,
                     font_size=8, color=RGBColor(0xBF, 0xDB, 0xFC),
                     align=PP_ALIGN.CENTER)


def bullet_block(slide, title, bullets, l, t, w, h,
                 title_color=C_INDIGO, bg=None, icon="•"):
    if bg:
        add_rect(slide, l, t, w, h, bg)
    add_text_box(slide, title, l+0.15, t+0.1, w-0.3, 0.38,
                 font_size=12, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(
        Inches(l+0.15), Inches(t+0.48), Inches(w-0.3), Inches(h-0.6))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon}  {b}"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = C_DARK


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# Full gradient-like background
add_rect(slide, 0, 0, 13.33, 7.5, C_INDIGO2)
add_rect(slide, 0, 5.2, 13.33, 2.3, RGBColor(0x06, 0xB6, 0xD4))  # cyan bottom band

# Decorative circles
for cx, cy, r, col in [
    (11.5, 1.2, 2.8, RGBColor(0x31, 0x2E, 0x90)),
    (12.5, 5.5, 2.0, RGBColor(0x06, 0x8F, 0xA8)),
]:
    sh = slide.shapes.add_shape(9, Inches(cx-r/2), Inches(cy-r/2),
                                 Inches(r), Inches(r))
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    sh.line.fill.background()

add_text_box(slide, "🏨  HOTEL BOOKING ANALYSIS", 0.7, 1.6, 11.9, 1.2,
             font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Comprehensive Business Intelligence Report",
             0.7, 2.85, 11.9, 0.7,
             font_size=20, color=C_CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide, "30,000 Hotel Bookings  |  10 US Cities  |  2024–2025",
             0.7, 3.55, 11.9, 0.5,
             font_size=13, color=RGBColor(0xBF, 0xDB, 0xFC), align=PP_ALIGN.CENTER)
add_text_box(slide, "Key Observations  •  Root Cause Analysis  •  Business Recommendations",
             0.7, 5.4, 11.9, 0.55,
             font_size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_text_box(slide, "Analysis Date: June 2026", 0.7, 6.8, 11.9, 0.4,
             font_size=10, color=C_DARK, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – EXECUTIVE SUMMARY (KPI Cards)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Executive Summary",
                 "At-a-glance key performance indicators across 30,000 hotel bookings")

# KPI row 1
kpi_data = [
    ("Total Bookings",     "30,000",  "2024–2025",                C_INDIGO),
    ("Overall Cancel Rate","20.23%",  "6,070 cancelled bookings",  C_RED),
    ("Web Channel Value",  "$28,191", "Highest avg booking value", C_GREEN),
    ("Travel Agent Cancel","27.9%",   "Highest cancellation rate", C_AMBER),
]
for i, (label, val, sub, bg) in enumerate(kpi_data):
    add_kpi_card(slide, 0.25 + i*3.24, 1.55, 3.0, 1.25, label, val, sub, bg)

# KPI row 2
kpi_data2 = [
    ("Confirmed Bookings", "21,672",  "72.2% of all bookings",     C_GREEN),
    ("Avg Stay Length",    "~4 nights","Consistent across all tiers",C_CYAN),
    ("Peak Month",         "May",     "2,910 bookings",             C_INDIGO),
    ("Deluxe Cancel Rate", "16.0%",   "Lowest cancellation room",   RGBColor(0x06,0x95,0x72)),
]
for i, (label, val, sub, bg) in enumerate(kpi_data2):
    add_kpi_card(slide, 0.25 + i*3.24, 3.05, 3.0, 1.25, label, val, sub, bg)

# Summary text
add_rect(slide, 0.25, 4.55, 12.83, 2.7, C_WHITE)
add_text_box(slide,
    "This report analyses 30,000 hotel bookings across 10 major US cities (2024–2025), covering three booking channels, "
    "three room types, and hotels rated 2–5 stars. The dataset reveals a 20.23% overall cancellation rate, with Travel Agents "
    "driving the highest cancellation (27.9%) while the Web channel delivers the best booking quality (17.6% cancel, $28,191 avg value). "
    "Seasonal peaks occur in April–June (summer + weddings); August–September show elevated cancellations (10–11%). "
    "Five-star properties carry slightly higher cancellation risk (21.3%) due to price sensitivity, while Deluxe rooms show the "
    "strongest guest commitment (16.0%). Recommendations focus on deposit policies, direct-booking loyalty programmes, and "
    "dynamic pricing to improve profitability and retention.",
    0.4, 4.65, 12.5, 2.5, font_size=10, color=C_SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – KEY OBSERVATIONS: Trends 1 & 2
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Key Observations – Trends & Patterns (1/2)",
                 "Booking channel dominance, seasonal demand cycles, and their interaction with value")

# Trend 1 box – Channel performance
add_rect(slide, 0.25, 1.55, 6.0, 5.7, C_WHITE)
add_text_box(slide, "📡  Trend 1: Channel Performance Divergence",
             0.4, 1.65, 5.7, 0.4, font_size=13, bold=True, color=C_INDIGO)
add_text_box(slide,
    "The three booking channels show dramatically different economics:\n\n"
    "• Web Channel (50% share): Highest avg value ($28,191) + lowest cancellation (17.6%)\n"
    "  → Customers research thoroughly before committing.\n\n"
    "• Mobile App (40% share): Mid-range value ($21,351), 21.6% cancel rate.\n"
    "  → Impulse/deal-driven bookings via push notifications.\n\n"
    "• Travel Agents (10% share): $24,454 avg value, but 27.9% cancellation.\n"
    "  → Agents hold multiple options, leading to high attrition.\n\n"
    "Key takeaway: The Web channel is the most valuable segment — 37% lower\n"
    "cancellation than Travel Agents, 32% higher average booking value than Mobile.",
    0.4, 2.1, 5.7, 4.5, font_size=10, color=C_SLATE)

# Trend 2 box – Seasonal patterns
add_rect(slide, 6.58, 1.55, 6.5, 5.7, C_WHITE)
add_text_box(slide, "📅  Trend 2: Seasonal Demand Swings",
             6.73, 1.65, 6.2, 0.4, font_size=13, bold=True, color=C_INDIGO)
add_text_box(slide,
    "Monthly booking volumes and cancellation rates move inversely:\n\n"
    "• Peak bookings (Apr–Jun): 2,221–2,910/month. Cancellation rate is\n"
    "  remarkably low (0.4–0.6%) — guests are committed to planned trips.\n\n"
    "• Post-summer surge (Aug–Sep): Volumes drop to ~1,864–1,920 but\n"
    "  cancellation rates spike to 10–11% — last-minute plan collapses.\n\n"
    "• Winter months (Jan–Feb): Moderate volumes (1,742–1,959), low-mid\n"
    "  cancellation (7–8%) — post-holiday, deliberate bookings.\n\n"
    "• Average booking value is relatively stable ($24K–$27K year-round),\n"
    "  suggesting price is not the main driver of seasonal variation —\n"
    "  traveller intent and external factors (school calendars, weather)\n"
    "  are the primary levers.",
    6.73, 2.1, 6.2, 4.5, font_size=10, color=C_SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – KEY OBSERVATIONS: Trend 3 + Cancellation overview
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Key Observations – Trends & Patterns (2/2)",
                 "Room type commitment levels, star rating nuances, and cancellation behaviour")

# Trend 3 – Room type
add_rect(slide, 0.25, 1.55, 6.0, 2.7, C_WHITE)
add_text_box(slide, "🛏️  Trend 3: Room Type & Commitment Levels",
             0.4, 1.65, 5.7, 0.4, font_size=13, bold=True, color=C_INDIGO)
add_text_box(slide,
    "• Standard (55.2% of bookings): Cancellation rate 23.3% — budget guests\n"
    "  have more alternatives; lower switching cost drives higher attrition.\n\n"
    "• Deluxe (34.9%): Cancellation rate 16.0% — strongest commitment;\n"
    "  guests are willing to pay more and see through their plans.\n\n"
    "• Suite (9.9%): Cancellation rate 18.0% — premium but moderate.",
    0.4, 2.1, 5.7, 2.0, font_size=10, color=C_SLATE)

# Trend 4 – Star rating
add_rect(slide, 6.58, 1.55, 6.5, 2.7, C_WHITE)
add_text_box(slide, "⭐  Trend 4: Star Rating & Cancellation Nuance",
             6.73, 1.65, 6.2, 0.4, font_size=13, bold=True, color=C_INDIGO)
add_text_box(slide,
    "• 2-Star: 19.8% cancel — lowest among all ratings; budget guests less\n"
    "  likely to have premium alternatives to switch to.\n\n"
    "• 3-Star: 20.2%  |  4-Star: 20.0% — near-identical, mid-market stable.\n\n"
    "• 5-Star: 21.3% — highest; luxury price sensitivity and broader high-end\n"
    "  alternatives increase cancellation risk.\n\n"
    "• Average stay (~4 nights) is consistent across all star ratings.",
    6.73, 2.1, 6.2, 2.0, font_size=10, color=C_SLATE)

# Cancellation overview table (bottom half)
add_rect(slide, 0.25, 4.5, 12.83, 2.75, C_WHITE)
add_text_box(slide, "📊  Cancellation Behaviour at a Glance",
             0.4, 4.58, 12.5, 0.4, font_size=13, bold=True, color=C_INDIGO)

rows = [
    ("Segment",          "Volume",  "Cancel Rate", "Avg Value", "Key Driver"),
    ("Web Channel",      "15,001",  "17.6%",       "$28,191",   "High intent, direct research"),
    ("Mobile App",       "12,009",  "21.6%",       "$21,351",   "Impulse booking, easy cancellation"),
    ("Travel Agent",     "2,990",   "27.9%",       "$24,454",   "Speculative holds, no deposit"),
    ("Standard Room",    "16,552",  "23.3%",       "$25,147",   "Low switching cost for budget guests"),
    ("Deluxe Room",      "10,478",  "16.0%",       "$25,005",   "Higher commitment, curated choice"),
    ("5-Star Property",  "4,511",   "21.3%",       "$25,117",   "Price sensitivity, luxury alternatives"),
]

col_widths = [2.5, 1.3, 1.3, 1.3, 6.1]
col_x = [0.4, 2.95, 4.3, 5.65, 7.0]
header_bg = C_INDIGO
row_bgs   = [C_WHITE, RGBColor(0xEF, 0xF6, 0xFF)]

for r_idx, row in enumerate(rows):
    bg = header_bg if r_idx == 0 else row_bgs[r_idx % 2]
    txt_color = C_WHITE if r_idx == 0 else C_DARK
    bold = r_idx == 0
    row_t = 5.0 + r_idx * 0.32
    add_rect(slide, 0.35, row_t, 12.63, 0.32, bg)
    for c_idx, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
        add_text_box(slide, cell, cx, row_t + 0.04, cw - 0.05, 0.28,
                     font_size=8.5, bold=bold, color=txt_color)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – VISUALIZATIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Data Visualizations",
                 "Cancellation rates, booking values, seasonal trends, and stay-length distribution")

if os.path.exists(VIZ_PATH):
    slide.shapes.add_picture(VIZ_PATH, Inches(0.15), Inches(1.45),
                              Inches(13.03), Inches(5.9))
else:
    add_text_box(slide, "⚠️  Run hotel_booking_analysis.py first to generate visualizations.png",
                 0.5, 3.5, 12.33, 0.6, font_size=14, color=C_RED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – ROOT CAUSE ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Root Cause Analysis",
                 "Why do cancellations, channel gaps, and seasonal patterns occur?")

# 3 columns
cols = [
    {
        "title"  : "🚫  Cancellation Drivers",
        "color"  : C_RED,
        "items"  : [
            "Travel Agents (27.9%): Speculative multi-property holds with no deposit risk; "
            "corporate flexible cancellation policies incentivise last-minute changes.",
            "Standard Rooms (23.3%): Budget guests face lower switching costs — identical "
            "or cheaper options are always available; no loyalty lock-in.",
            "5-Star Properties (21.3%): High absolute price amplifies sensitivity to any "
            "change in travel plans; abundance of luxury competitors in top US cities.",
            "Aug–Sep spike (10–11%): Post-summer plan collapses driven by school calendars, "
            "budget exhaustion, and weather uncertainty.",
        ]
    },
    {
        "title"  : "📡  Channel Performance Gap",
        "color"  : C_INDIGO,
        "items"  : [
            "Web superiority: Direct booking = high intentionality. Customers compare, "
            "research, and decide deliberately — 'booker's commitment' is highest.",
            "Mobile App underperformance: Push-notification flash deals attract "
            "price-sensitive, impulsive bookers who cancel when better prices emerge.",
            "Travel Agent structurally weak: Agents are paid per booking, not per stay. "
            "No financial penalty for cancellations means zero deterrent to over-booking.",
            "Avg markup is similar across channels (~$6,930–$6,985), so revenue "
            "loss from Travel Agent cancellations is not offset by higher margins.",
        ]
    },
    {
        "title"  : "📅  Seasonal & Temporal Factors",
        "color"  : C_GREEN,
        "items"  : [
            "Apr–Jun peak: Summer vacation planning + US wedding season creates committed "
            "forward bookings; cancellation rate drops to <1% — these guests are locked in.",
            "Aug–Sep trough: School year restart, family budget exhaustion, and weather "
            "shifts cause last-minute trip abandonment — a structural, recurring pattern.",
            "Stable stay length (~4 nights) across ALL star ratings suggests trip purpose "
            "(not property tier) drives duration — segment marketing accordingly.",
            "Avg booking value stable year-round ($24K–$27K): Seasonality affects volume "
            "and commitment, not willingness-to-pay per booking.",
        ]
    },
]

for i, col in enumerate(cols):
    lx = 0.25 + i * 4.35
    add_rect(slide, lx, 1.55, 4.1, 5.7, C_WHITE)
    add_text_box(slide, col["title"], lx+0.15, 1.65, 3.8, 0.42,
                 font_size=12, bold=True, color=col["color"])
    txBox = slide.shapes.add_textbox(
        Inches(lx+0.15), Inches(2.12), Inches(3.8), Inches(5.0))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(col["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = C_SLATE

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – BUSINESS RECOMMENDATIONS: Cancellations
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Business Recommendations – Reduce Cancellations",
                 "Targeted policies to cut the 20.23% cancellation rate by 15–20%")

rec_blocks = [
    {
        "icon": "🏢", "title": "Travel Agent Channel (Priority 1 – Immediate)",
        "color": C_RED,
        "items": [
            "Require 25% non-refundable deposit for all agent-originated bookings.",
            "Implement agent performance scorecards — completion rate < 75% triggers commission tier demotion.",
            "'Verified Booking' protocol: agents must provide customer contact for confirmation call within 48 hrs.",
            "Offer agents a 'completion bonus' ($50/booking) to incentivise follow-through over cancellations.",
            "Expected impact: 8–12% reduction in Travel Agent cancellations.",
        ]
    },
    {
        "icon": "🛏️", "title": "Standard Room Segment (Priority 2 – 4 Weeks)",
        "color": C_AMBER,
        "items": [
            "'Flex-Date Credits': instead of cancelling, guests receive a 6-month credit — reduces refunds.",
            "Complimentary room upgrade (Standard → Deluxe) when guest accepts date change instead of cancellation.",
            "Trip insurance upsell at checkout: 'Cancel for Any Reason' at 3% of booking value.",
            "Send 48-hr pre-arrival reminder with value-add (restaurant coupon) to reinforce commitment.",
            "Expected impact: 5–7% reduction in Standard room cancellations.",
        ]
    },
    {
        "icon": "📅", "title": "Seasonal Intervention (Priority 3 – Before Aug/Sep)",
        "color": C_INDIGO,
        "items": [
            "For bookings with check-in in Aug–Sep, offer early 'lock-in incentive' (free breakfast if not cancelled).",
            "Stricter cancellation windows during Aug–Sep: non-refundable 14 days before check-in vs 3 days year-round.",
            "Mobile App: Proactive 'Trip Confirmation' push at T-30 days to confirm travel plans.",
            "Cross-sell travel insurance via Mobile App channel specifically for high-cancel months.",
        ]
    },
]

for i, rec in enumerate(rec_blocks):
    ty = 1.6 + i * 1.9
    add_rect(slide, 0.25, ty, 12.83, 1.8, C_WHITE)
    add_text_box(slide, f"{rec['icon']}  {rec['title']}",
                 0.4, ty+0.1, 12.5, 0.38, font_size=12, bold=True, color=rec["color"])
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(ty+0.5), Inches(12.5), Inches(1.25))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(rec["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"✔  {item}"
        for run in p.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = C_SLATE

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – BUSINESS RECOMMENDATIONS: Profitability & Repeat Bookings
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Business Recommendations – Profitability & Repeat Bookings",
                 "Strategies to increase direct revenue, customer lifetime value, and retention")

left_items = [
    ("📲  Direct Booking Loyalty Programme", C_INDIGO, [
        "Launch 'Book Direct Rewards': 10% instant discount + loyalty points for Web channel bookings.",
        "Points earn: 1 point per $10 spent; 500 points = 1 free night (drives repeat stays).",
        "Exclusive perks for direct bookers: late checkout, free WiFi, F&B credit ($25).",
        "Run 'Book Direct' awareness campaign across email + social media targeting past guests.",
    ]),
    ("🔄  Post-Stay Retention", C_GREEN, [
        "Personalised rebooking offer within 24 hours of checkout: 15% off same property.",
        "Birthday/anniversary recognition: room upgrade offer 30 days before special dates.",
        "Milestone loyalty rewards: 3rd stay (room upgrade), 5th stay (free night), 10th (VIP status).",
        "NPS survey post-stay; high-score guests get referral programme invitation.",
    ]),
]

right_items = [
    ("📈  Revenue Per Booking Optimisation", C_AMBER, [
        "Upsell Deluxe/Suite at check-in for Standard room guests (targeted via Mobile App).",
        "Bundle packages: Room + Breakfast ($15 add-on), Room + Airport Transfer, Room + Spa.",
        "Minimum-stay upsell: 'Book 3 nights, save 12%' for stays that would otherwise be 1–2 nights.",
        "Target repeat guests with higher-tier room offers based on their last stay preferences.",
    ]),
    ("🏢  Channel Revenue Mix", C_CYAN, [
        "Shift 5% of Travel Agent share to Web channel (saves cancellation losses + increases avg value).",
        "Reduce OTA dependency by investing in Web channel SEO and PPC campaigns.",
        "For Mobile App: shift from deal-focused to experience-focused marketing to raise avg value.",
        "Track channel contribution margin, not just revenue — factor in cancellation cost.",
    ]),
]

for i, (title, color, items) in enumerate(left_items):
    ty = 1.6 + i * 2.8
    add_rect(slide, 0.25, ty, 6.2, 2.65, C_WHITE)
    add_text_box(slide, title, 0.4, ty+0.1, 5.9, 0.38, font_size=11, bold=True, color=color)
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(ty+0.5), Inches(5.9), Inches(2.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"✔  {item}"
        for run in p.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = C_SLATE

for i, (title, color, items) in enumerate(right_items):
    ty = 1.6 + i * 2.8
    add_rect(slide, 6.88, ty, 6.2, 2.65, C_WHITE)
    add_text_box(slide, title, 7.03, ty+0.1, 5.9, 0.38, font_size=11, bold=True, color=color)
    txBox = slide.shapes.add_textbox(Inches(7.03), Inches(ty+0.5), Inches(5.9), Inches(2.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"✔  {item}"
        for run in p.runs:
            run.font.size = Pt(9.5)
            run.font.color.rgb = C_SLATE

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – BUSINESS RECOMMENDATIONS: Pricing & Channel Strategy
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Business Recommendations – Pricing, Promotions & Channel Strategy",
                 "Optimise revenue yield through dynamic pricing and segment-specific promotions")

# Channel strategy cards (top row)
ch_cards = [
    ("🌐  Web Channel", C_INDIGO,
     "Premium positioning & exclusivity.\n"
     "• Price parity + exclusive perks (late checkout, F&B credit)\n"
     "• 'Best Rate Guarantee' to prevent OTA leakage\n"
     "• Enhanced UX: one-click rebooking for returning guests\n"
     "• Retargeting ads for cart-abandoned users"),
    ("📱  Mobile App", C_CYAN,
     "Flash sales & last-minute deals.\n"
     "• 18-hour flash sales via push notifications (mobile-exclusive)\n"
     "• Progressive discount: book within 1hr → extra 5% off\n"
     "• Gamified check-in streaks and app-only loyalty points\n"
     "• AI-driven 'Best Deals Near You' personalisation"),
    ("🤝  Travel Agents", C_AMBER,
     "Volume-based commission reform.\n"
     "• Tiered commissions: 8% (base), 12% (50+/month), 15% (100+/month)\n"
     "• Commission tied to stay completion, not just booking\n"
     "• Quarterly agent performance review and ranking\n"
     "• Co-marketing budget for top-performing agents"),
]
for i, (title, color, body) in enumerate(ch_cards):
    lx = 0.25 + i * 4.35
    add_rect(slide, lx, 1.55, 4.1, 2.55, C_WHITE)
    add_text_box(slide, title, lx+0.15, 1.65, 3.8, 0.38, font_size=12, bold=True, color=color)
    add_text_box(slide, body, lx+0.15, 2.1, 3.8, 1.9, font_size=9.5, color=C_SLATE)

# Seasonal promotions
add_rect(slide, 0.25, 4.3, 6.2, 2.95, C_WHITE)
add_text_box(slide, "🎯  Seasonal Promotions", 0.4, 4.4, 5.9, 0.38,
             font_size=12, bold=True, color=C_RED)
promo_text = (
    "• Apr–Jun (Peak): 'Early Bird Summer' — 15% off for 60+ day advance bookings.\n"
    "  Minimum 2-night stay requirement to maximise occupancy yield.\n\n"
    "• Aug–Sep (High Cancel): 'Trip Protection Bundle' — book with trip insurance;\n"
    "  offer 'rebook, don't cancel' credits valid 9 months.\n\n"
    "• Feb (Slow): 'Valentines & Beyond' — couples packages, extended stay promo\n"
    "  (stay 4 nights, 5th free) to stimulate demand.\n\n"
    "• Oct–Dec (Steady): 'Holiday Countdown' early access for loyalty members."
)
add_text_box(slide, promo_text, 0.4, 4.85, 5.9, 2.3, font_size=9.5, color=C_SLATE)

# Dynamic pricing
add_rect(slide, 6.88, 4.3, 6.2, 2.95, C_WHITE)
add_text_box(slide, "⚙️  Dynamic Pricing Engine", 7.03, 4.4, 5.9, 0.38,
             font_size=12, bold=True, color=C_GREEN)
pricing_text = (
    "• Real-time rate adjustment based on 30/60/90-day occupancy forecasts.\n\n"
    "• Competitor rate benchmarking: auto-match or undercut by ≤3% when\n"
    "  occupancy is below 70% with 30+ days to check-in.\n\n"
    "• Last-minute pricing (< 7 days): Mobile App exclusive discounts (up to 20%)\n"
    "  to fill unsold inventory without devaluing brand on main channels.\n\n"
    "• Peak period floor prices: Set minimum rates during Apr–Jun to protect\n"
    "  revenue when demand is inelastic.\n\n"
    "• Room-type-specific pricing: Deluxe/Suite premium to reflect lower\n"
    "  cancellation risk and higher perceived value."
)
add_text_box(slide, pricing_text, 7.03, 4.85, 5.9, 2.3, font_size=9.5, color=C_SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – IMPLEMENTATION ROADMAP & EXPECTED IMPACT
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)
add_slide_header(slide, "Implementation Roadmap & Expected Impact",
                 "Phased execution plan with measurable KPI targets")

phases = [
    ("Phase 1", "Immediate\n(0–4 Weeks)", C_RED,
     ["Agent deposit requirement (25% non-refundable)",
      "Post-stay email campaign launch",
      "'Flex-Date Credits' for Standard room cancellations",
      "Aug/Sep trip insurance upsell activation"]),
    ("Phase 2", "Short-term\n(1–3 Months)", C_AMBER,
     ["Book Direct Rewards loyalty programme launch",
      "Mobile App flash sale system (18-hr windows)",
      "Agent performance scorecard system",
      "Early-bird summer packages (Apr–Jun)"]),
    ("Phase 3", "Mid-term\n(3–6 Months)", C_INDIGO,
     ["Dynamic pricing engine deployment",
      "Bundle package creation (room + add-ons)",
      "Mobile App personalisation (AI recommendations)",
      "NPS + referral programme rollout"]),
    ("Phase 4", "Long-term\n(6–12 Months)", C_GREEN,
     ["Full loyalty milestone reward system",
      "Agent commission tier restructure complete",
      "Competitor benchmarking + automated repricing",
      "Cross-channel analytics dashboard"]),
]

for i, (phase, timeline, color, actions) in enumerate(phases):
    lx = 0.25 + i * 3.26
    add_rect(slide, lx, 1.55, 3.0, 0.6, color)
    add_text_box(slide, phase, lx+0.1, 1.58, 2.8, 0.28,
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, timeline, lx+0.1, 1.85, 2.8, 0.28,
                 font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, lx, 2.18, 3.0, 3.1, C_WHITE)
    for j, action in enumerate(actions):
        add_text_box(slide, f"✔  {action}", lx+0.12, 2.28 + j*0.72, 2.76, 0.65,
                     font_size=9.5, color=C_SLATE)

# Impact KPIs (bottom)
add_rect(slide, 0.25, 5.5, 12.83, 1.75, C_INDIGO2)
add_text_box(slide, "EXPECTED IMPACT AFTER FULL IMPLEMENTATION",
             0.4, 5.56, 12.5, 0.38, font_size=13, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

impacts = [
    ("15–20%", "Reduction in\nOverall Cancellation Rate"),
    ("10–15%", "Increase in\nDirect (Web) Bookings"),
    ("8–12%", "Improvement in\nAvg Booking Value"),
    ("25%", "Increase in\nRepeat Customer Rate"),
]
for i, (val, label) in enumerate(impacts):
    lx = 0.6 + i * 3.1
    add_text_box(slide, val, lx, 5.97, 2.6, 0.5,
                 font_size=22, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, lx, 6.48, 2.6, 0.6,
                 font_size=9, color=RGBColor(0xBF, 0xDB, 0xFC), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – CONCLUSION
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_INDIGO2)
add_rect(slide, 0, 5.5, 13.33, 2.0, C_CYAN)

add_text_box(slide, "Conclusion & Next Steps", 0.7, 1.0, 11.9, 0.8,
             font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

conclusion = (
    "This analysis of 30,000 hotel bookings across 10 US cities reveals a clear strategic opportunity:\n\n"
    "The 20.23% cancellation rate — costing millions in lost revenue and operational disruption — is "
    "largely driven by structural weaknesses in the Travel Agent channel (27.9%), low-commitment Standard "
    "room bookings (23.3%), and a recurring August–September cancellation spike (10–11%).\n\n"
    "The Web channel, despite commanding only 50% of bookings, is the gold standard: lowest cancellation (17.6%), "
    "highest booking value ($28,191), and strongest guest commitment. Growing this channel is the single highest-ROI action.\n\n"
    "The four-phase roadmap — from immediate agent deposit policies to long-term dynamic pricing — provides a "
    "structured path to 15–20% cancellation reduction, 10–15% direct booking growth, and a 25% repeat customer rate."
)
add_text_box(slide, conclusion, 0.7, 1.95, 11.9, 3.3, font_size=11, color=C_WHITE)

add_text_box(slide, "Immediate Next Step: Implement Travel Agent deposit requirement and Book Direct Rewards programme",
             0.7, 5.6, 11.9, 0.5, font_size=12, bold=True,
             color=C_DARK, align=PP_ALIGN.CENTER)
add_text_box(slide, "Hotel Booking Analysis  |  30,000 Bookings  |  June 2026",
             0.7, 6.9, 11.9, 0.4, font_size=10, color=C_DARK, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(REPORT_PATH)
print(f"  [OK] PowerPoint saved: {REPORT_PATH}")
